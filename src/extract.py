"""Extract text from source artifacts (read-only) into pilot/corpus/.

Source artifacts are NEVER modified or deleted. We only read from them and
write extracted text to a separate corpus directory.
"""
from __future__ import annotations
import os
import re
import sys
from pathlib import Path

import fitz  # PyMuPDF
from pptx import Presentation
from bs4 import BeautifulSoup

SRC = Path(os.environ.get("DOC_SOURCE_DIR", Path(__file__).resolve().parent.parent / "source_docs"))
DST = Path(__file__).resolve().parent.parent / "corpus"

SKIP_EXT = {".png", ".jpg", ".jpeg", ".gif", ".gslides", ".DS_Store", ".js"}
SKIP_NAMES = {".DS_Store", "build_pptx.js"}


def safe_name(name: str) -> str:
    return re.sub(r"[^\w\-.]+", "_", name)


def extract_pdf(path: Path) -> str:
    parts = []
    with fitz.open(path) as doc:
        for i, page in enumerate(doc):
            t = page.get_text("text") or ""
            t = t.strip()
            if t:
                parts.append(f"\n\n--- page {i+1} ---\n{t}")
    return "\n".join(parts).strip()


def extract_pptx(path: Path) -> str:
    parts = []
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides):
        chunks = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text or "" for run in para.runs).strip()
                    if line:
                        chunks.append(line)
            elif shape.shape_type == 19:  # table
                pass
        if chunks:
            parts.append(f"\n\n--- slide {i+1} ---\n" + "\n".join(chunks))
    return "\n".join(parts).strip()


def extract_html(path: Path) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        html = f.read()
    soup = BeautifulSoup(html, "lxml")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    text = soup.get_text(separator="\n")
    lines = [ln.strip() for ln in text.splitlines()]
    lines = [ln for ln in lines if ln]
    return "\n".join(lines)


def extract_txt(path: Path) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def main() -> int:
    DST.mkdir(parents=True, exist_ok=True)
    manifest = []
    failures = []
    for p in sorted(SRC.iterdir()):
        if not p.is_file():
            continue
        if p.name in SKIP_NAMES or p.suffix.lower() in SKIP_EXT:
            continue
        ext = p.suffix.lower()
        try:
            if ext == ".pdf":
                text = extract_pdf(p)
            elif ext == ".pptx":
                text = extract_pptx(p)
            elif ext in (".html", ".htm"):
                text = extract_html(p)
            elif ext in (".txt", ".md"):
                text = extract_txt(p)
            else:
                continue
        except Exception as e:
            failures.append((p.name, repr(e)))
            continue
        if not text or len(text) < 50:
            failures.append((p.name, "empty or too short"))
            continue
        out_name = safe_name(p.stem) + ".txt"
        out_path = DST / out_name
        with open(out_path, "w", encoding="utf-8") as f:
            f.write(f"# SOURCE: {p.name}\n")
            f.write(f"# BYTES: {p.stat().st_size}\n\n")
            f.write(text)
        manifest.append((p.name, out_name, len(text)))
        print(f"  ok   {p.name:60s} -> {out_name}  ({len(text):,} chars)")

    print(f"\nExtracted {len(manifest)} files into {DST}")
    if failures:
        print(f"\n{len(failures)} failures:")
        for name, err in failures:
            print(f"  fail {name}: {err}")
    return 0 if manifest else 1


if __name__ == "__main__":
    sys.exit(main())
